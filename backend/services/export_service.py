import json
from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptInches
from io import BytesIO


def export_to_word(activities: dict) -> BytesIO:
    doc = Document()
    doc.add_heading(f"الأنشطة التعليمية - {activities.get('lesson', '')}", 0)

    doc.add_heading("أسئلة الاختيار المتعدد", level=1)
    for i, q in enumerate(activities.get("mcq", []), 1):
        doc.add_paragraph(f"{i}. {q['question']}")
        for choice in q.get("choices", []):
            doc.add_paragraph(f"   • {choice}", style="List Bullet")
        p = doc.add_paragraph(f"الإجابة: {q['answer']}")
        p.runs[0].bold = True

    doc.add_heading("صح وخطأ", level=1)
    for i, q in enumerate(activities.get("true_false", []), 1):
        doc.add_paragraph(f"{i}. {q['question']}")
        p = doc.add_paragraph(f"الإجابة: {'صحيح' if q['answer'] else 'خطأ'}")
        p.runs[0].bold = True

    doc.add_heading("البطاقات التعليمية", level=1)
    for i, card in enumerate(activities.get("flashcards", []), 1):
        doc.add_paragraph(f"{i}. السؤال: {card['front']}")
        doc.add_paragraph(f"   الجواب: {card['back']}")

    doc.add_heading("لعبة المطابقة", level=1)
    for pair in activities.get("matching", []):
        doc.add_paragraph(f"{pair['left']}  ←  {pair['right']}")

    doc.add_heading("خطة الدرس", level=1)
    plan = activities.get("lesson_plan", {})
    if plan:
        doc.add_paragraph(f"العنوان: {plan.get('title', '')}")
        doc.add_paragraph(f"المرحلة: {plan.get('grade', '')}")
        for obj in plan.get("objectives", []):
            doc.add_paragraph(f"• {obj}", style="List Bullet")

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def export_to_pptx(activities: dict) -> BytesIO:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"الأنشطة التعليمية - {activities.get('lesson', '')}"

    for i, q in enumerate(activities.get("mcq", []), 1):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"سؤال {i}: {q['question']}"
        content = "\n".join(f"• {c}" for c in q.get("choices", []))
        slide.placeholders[1].text = content + f"\n\nالإجابة: {q['answer']}"

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def export_to_pdf_text(activities: dict) -> str:
    lines = []
    lines.append(f"الأنشطة التعليمية - {activities.get('lesson', '')}")
    lines.append("=" * 50)

    lines.append("\nأسئلة الاختيار المتعدد:")
    for i, q in enumerate(activities.get("mcq", []), 1):
        lines.append(f"{i}. {q['question']}")
        for choice in q.get("choices", []):
            lines.append(f"   - {choice}")
        lines.append(f"   الإجابة: {q['answer']}")

    lines.append("\nصح وخطأ:")
    for i, q in enumerate(activities.get("true_false", []), 1):
        lines.append(f"{i}. {q['question']} ({'صحيح' if q['answer'] else 'خطأ'})")

    lines.append("\nالبطاقات التعليمية:")
    for i, card in enumerate(activities.get("flashcards", []), 1):
        lines.append(f"{i}. {card['front']} -> {card['back']}")

    lines.append("\nالمطابقة:")
    for pair in activities.get("matching", []):
        lines.append(f"  {pair['left']} ←→ {pair['right']}")

    return "\n".join(lines)


def format_for_moodle(activities: dict) -> str:
    xml_parts = ['<?xml version="1.0" encoding="UTF-8"?>\n<quiz>']
    for q in activities.get("mcq", []):
        xml_parts.append("  <question type='multichoice'>")
        xml_parts.append(f"    <name><text>{q['question'][:50]}</text></name>")
        xml_parts.append(f"    <questiontext format='html'><text>{q['question']}</text></questiontext>")
        for j, choice in enumerate(q.get("choices", [])):
            fraction = "100" if choice == q["answer"] else "0"
            xml_parts.append(f"    <answer fraction='{fraction}'><text>{choice}</text></answer>")
        xml_parts.append("  </question>")
    xml_parts.append("</quiz>")
    return "\n".join(xml_parts)


def format_for_google_forms(activities: dict) -> list:
    items = []
    for q in activities.get("mcq", []):
        items.append({
            "title": q["question"],
            "type": "MULTIPLE_CHOICE",
            "options": q.get("choices", []),
            "answer": q["answer"]
        })
    return items
